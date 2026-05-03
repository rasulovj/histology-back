import os
import json
import asyncio
import gc  # Добавлен сборщик мусора для очистки RAM
from config import DB_PATH
import docx
from pptx import Presentation
from pypdf import PdfReader

KB_DIR = DB_PATH
INDEX_FILE = os.path.join(KB_DIR, "library_index.json")
KNOWLEDGE_MAP_FILE = os.path.join(KB_DIR, "knowledge_map.json")
CACHE_DIR = os.path.join(KB_DIR, "text_cache") 
MAX_CHUNKS = 3000
MAX_TOTAL_TEXT_CHARS = 20_000_000
LARGE_FILE_SKIP_KMAP_BYTES = 120 * 1024 * 1024
MAX_WORDS_PER_CHUNK = 800

def load_json(filepath):
    if os.path.exists(filepath):
        try:
            with open(filepath, 'r', encoding='utf-8') as f:
                return json.load(f)
        except Exception:
            return {}
    return {}

def save_json(filepath, data):
    with open(filepath, 'w', encoding='utf-8') as f:
        json.dump(data, f, ensure_ascii=False, indent=4)

def extract_text_in_chunks(file_path):
    ext = file_path.lower().split('.')[-1]
    chunks = []
    total_chars = 0
    
    try:
        if ext == 'pdf':
            reader = PdfReader(file_path)
            for p in reader.pages:
                text = p.extract_text()
                normalized = text if text else ""
                chunks.append(normalized)
                total_chars += len(normalized)
                if len(chunks) >= MAX_CHUNKS or total_chars >= MAX_TOTAL_TEXT_CHARS:
                    break
                
        elif ext == 'pptx':
            prs = Presentation(file_path)
            for slide in prs.slides:
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)
                normalized = "\n".join(slide_text)
                chunks.append(normalized)
                total_chars += len(normalized)
                if len(chunks) >= MAX_CHUNKS or total_chars >= MAX_TOTAL_TEXT_CHARS:
                    break
                
        elif ext == 'docx':
            doc = docx.Document(file_path)
            chunk_size = 2000
            buffer = []
            buffer_len = 0
            for para in doc.paragraphs:
                para_text = para.text or ""
                if not para_text:
                    continue
                buffer.append(para_text)
                buffer_len += len(para_text) + 1
                if buffer_len >= chunk_size:
                    normalized = "\n".join(buffer)
                    chunks.append(normalized)
                    total_chars += len(normalized)
                    buffer = []
                    buffer_len = 0
                    if len(chunks) >= MAX_CHUNKS or total_chars >= MAX_TOTAL_TEXT_CHARS:
                        break
            if buffer and len(chunks) < MAX_CHUNKS and total_chars < MAX_TOTAL_TEXT_CHARS:
                normalized = "\n".join(buffer)
                chunks.append(normalized)
            
    except Exception as e:
        print(f"Ошибка чтения файла {file_path}: {e}")
        
    return chunks

async def index_document(file_path, doc_id, file_name):
    from services.ai_service import classify_book_topic
    
    print(f"🔄 Индексирую документ: {file_name}")
    
    chunks = await asyncio.to_thread(extract_text_in_chunks, file_path)
    
    cache_path = os.path.join(CACHE_DIR, f"{doc_id}.json")
    await asyncio.to_thread(save_json, cache_path, chunks)
    
    sample_text = "\n".join(chunks[:15])
    
    # Здесь вызывается твой Gemini API для классификации!
    category = await classify_book_topic(sample_text)

    file_size = 0
    try:
        file_size = os.path.getsize(file_path)
    except Exception:
        file_size = 0

    if file_size >= LARGE_FILE_SKIP_KMAP_BYTES:
        print(f"⚠️ Пропускаю heavy keyword-index для большого файла: {file_name} ({file_size} bytes)")
        del chunks
        gc.collect()
        return category
    
    k_map = await asyncio.to_thread(load_json, KNOWLEDGE_MAP_FILE)
    
    for page_num, text in enumerate(chunks):
        await asyncio.sleep(0)
        
        if not text:
            continue
            
        # Увеличили лимит до 5 букв, чтобы отсеять мусор и спасти RAM
        words = set(w.lower() for w in text.split() if 5 < len(w) < 40)
        if len(words) > MAX_WORDS_PER_CHUNK:
            words = set(list(words)[:MAX_WORDS_PER_CHUNK])
        
        for w in words:
            if w not in k_map:
                k_map[w] = []
                
            # МГНОВЕННАЯ ПРОВЕРКА вместо зависающего any()
            if not k_map[w] or not (k_map[w][-1]['b'] == doc_id and k_map[w][-1]['p'] == page_num):
                k_map[w].append({"b": doc_id, "p": page_num})
                
    await asyncio.to_thread(save_json, KNOWLEDGE_MAP_FILE, k_map)
    
    # ПРИНУДИТЕЛЬНАЯ ОЧИСТКА ПАМЯТИ СЕРВЕРА
    del k_map
    del chunks
    gc.collect()
    
    return category

async def sync_library():
    os.makedirs(KB_DIR, exist_ok=True)
    os.makedirs(CACHE_DIR, exist_ok=True) 
    
    index = await asyncio.to_thread(load_json, INDEX_FILE)
    allowed_exts = ('.pdf', '.docx', '.pptx')
    
    current_files = await asyncio.to_thread(os.listdir, KB_DIR)
    current_files = [f for f in current_files if f.lower().endswith(allowed_exts)]
    
    updated = False
    for i, file_name in enumerate(current_files):
        doc_id = f"doc_{i}"
        
        if doc_id not in index or index[doc_id]['title'] != file_name:
            file_path = os.path.join(KB_DIR, file_name)
            category = await index_document(file_path, doc_id, file_name)
            
            index[doc_id] = {
                "id": doc_id,
                "title": file_name,
                "path": file_path,
                "category": category
            }
            updated = True
            
    if updated:
        await asyncio.to_thread(save_json, INDEX_FILE, index)
        
    print("✅ Библиотека AI Study Assistant синхронизирована!")

async def get_library_catalog():
    return await asyncio.to_thread(load_json, INDEX_FILE)
